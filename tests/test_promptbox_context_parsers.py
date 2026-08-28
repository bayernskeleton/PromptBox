import sys

import pytest

from promptbox_mvp.context_parsers import parse_context_file


def test_json_success_is_normalized_and_bad_json_is_failed(tmp_path):
    good = tmp_path / "data.json"
    good.write_text('{"name":"数云"}', encoding="utf-8")
    bad = tmp_path / "bad.json"
    bad.write_text("{bad", encoding="utf-8")

    parsed_good = parse_context_file(good)
    parsed_bad = parse_context_file(bad)

    assert parsed_good.status == "success"
    assert '"name": "数云"' in parsed_good.text
    assert parsed_bad.status == "failed"
    assert parsed_bad.reason_code == "json_parse_failed"
    assert parsed_bad.text == ""


def test_text_file_uses_encoding_fallback(tmp_path):
    path = tmp_path / "说明.txt"
    path.write_bytes("中文内容".encode("gbk"))

    result = parse_context_file(path)

    assert result.status == "success"
    assert result.encoding == "gbk"
    assert result.text == "中文内容"


def test_csv_result_is_markdown_and_marks_row_truncation(tmp_path):
    path = tmp_path / "data.csv"
    path.write_text("名称,数值\n甲,1\n乙,2\n", encoding="utf-8")

    result = parse_context_file(path, csv_row_limit=1)

    assert result.status == "success"
    assert "| 名称 | 数值 |" in result.text
    assert result.truncated is True
    assert "前 1 行" in result.reason


def test_pptx_result_contains_slide_number_and_text(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    path = tmp_path / "brief.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "业务目标"
    presentation.save(path)

    result = parse_context_file(path)

    assert result.status == "success"
    assert "幻灯片：1" in result.text
    assert "业务目标" in result.text


def test_docx_extracts_paragraph_and_table_text(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("项目结论")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "值"
    document.save(path)

    result = parse_context_file(path)

    assert result.status == "success"
    assert "项目结论" in result.text
    assert "字段" in result.text
    assert "值" in result.text


def test_xlsx_extracts_sheet_name_and_cells(tmp_path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    path = tmp_path / "metrics.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "指标"
    sheet.append(["名称", "数值"])
    sheet.append(["转化率", 0.42])
    workbook.save(path)

    result = parse_context_file(path)

    assert result.status == "success"
    assert "工作表：指标" in result.text
    assert "转化率" in result.text
    assert "0.42" in result.text


def test_pdf_with_no_text_layer_is_explicit_failure(tmp_path):
    pytest.importorskip("pypdf")
    from pypdf import PdfWriter

    path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with path.open("wb") as handle:
        writer.write(handle)

    result = parse_context_file(path)

    assert result.status == "failed"
    assert result.reason_code == "empty_text_layer"
    assert result.text == ""


def test_unsupported_old_office_format_is_structured_failure(tmp_path):
    path = tmp_path / "legacy.doc"
    path.write_bytes(b"old office")

    result = parse_context_file(path)

    assert result.status == "failed"
    assert result.reason_code == "unsupported_format"
    assert ".docx" in result.reason


def test_missing_optional_dependency_is_structured_failure(tmp_path, monkeypatch):
    path = tmp_path / "report.docx"
    path.write_bytes(b"not a real docx")
    monkeypatch.setitem(sys.modules, "docx", None)

    result = parse_context_file(path)

    assert result.status == "failed"
    assert result.reason_code == "dependency_missing"
    assert "python-docx" in result.reason
