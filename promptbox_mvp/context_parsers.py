"""Lazy, structured parsers for files used as business context."""

import csv
import io
import json
from dataclasses import dataclass
from pathlib import Path

from .context_loader import (
    DEFAULT_CSV_ROW_LIMIT,
    DEFAULT_MAX_BYTES,
    _csv_to_markdown,
    read_text_with_fallback,
)
from .context_manifest import SUPPORTED_SUFFIXES

TEXT_SUFFIXES = SUPPORTED_SUFFIXES - {".json", ".csv", ".pdf", ".docx", ".xlsx", ".pptx"}
OPTIONAL_PACKAGES = {
    ".pdf": "pypdf",
    ".docx": "python-docx",
    ".xlsx": "openpyxl",
    ".pptx": "python-pptx",
}


@dataclass(frozen=True)
class ParserResult:
    status: str
    text: str
    chars: int
    encoding: str | None = None
    truncated: bool = False
    reason_code: str | None = None
    reason: str = ""
    parser_version: str = "context-parser-v1"


def _success(text: str, *, encoding: str | None = None, truncated: bool = False, reason: str = "") -> ParserResult:
    return ParserResult("success", text, len(text), encoding, truncated, None, reason)


def _failure(code: str, reason: str) -> ParserResult:
    return ParserResult("failed", "", 0, None, False, code, reason)


def _parse_json(path: Path) -> ParserResult:
    try:
        text, encoding = read_text_with_fallback(path)
        parsed = json.loads(text)
    except (UnicodeError, ValueError, json.JSONDecodeError) as exc:
        return _failure("json_parse_failed", f"JSON 解析失败：{exc}")
    return _success(json.dumps(parsed, indent=2, ensure_ascii=False), encoding=encoding)


def _parse_csv(path: Path, row_limit: int) -> ParserResult:
    try:
        text, encoding = read_text_with_fallback(path)
        converted, truncated, total = _csv_to_markdown(text, row_limit)
    except (OSError, UnicodeError, csv.Error, ValueError) as exc:
        return _failure("document_read_failed", f"CSV 读取失败：{exc}")
    if not converted:
        return _failure("empty_content", "CSV 文件没有可读取内容")
    reason = f"CSV 共 {total} 行数据，已截断显示前 {row_limit} 行。" if truncated else ""
    return _success(converted, encoding=encoding, truncated=truncated, reason=reason)


def _parse_pdf(path: Path) -> ParserResult:
    try:
        from pypdf import PdfReader
    except (ImportError, ModuleNotFoundError):
        return _failure("dependency_missing", "缺少 pypdf，请安装 pypdf 后重试，或复制正文到剪贴板")
    try:
        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(f"[页码：{index}]\n{text}")
    except Exception as exc:
        return _failure("document_read_failed", f"PDF 读取失败：{exc}")
    if not pages:
        return _failure("empty_text_layer", "PDF 没有可复制文本层；当前版本不做 OCR")
    return _success("\n\n".join(pages))


def _parse_docx(path: Path) -> ParserResult:
    try:
        from docx import Document
    except (ImportError, ModuleNotFoundError):
        return _failure("dependency_missing", "缺少 python-docx，请安装 python-docx 后重试，或复制正文到剪贴板")
    try:
        document = Document(str(path))
        blocks = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            rows = [" | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells) for row in table.rows]
            if rows:
                blocks.append("[表格开始]\n" + "\n".join(rows) + "\n[表格结束]")
    except Exception as exc:
        return _failure("document_read_failed", f"DOCX 读取失败：{exc}；可复制正文到剪贴板")
    text = "\n\n".join(blocks)
    return _success(text) if text else _failure("empty_content", "DOCX 没有可读取内容")


def _parse_xlsx(path: Path, row_limit: int) -> ParserResult:
    try:
        from openpyxl import load_workbook
    except (ImportError, ModuleNotFoundError):
        return _failure("dependency_missing", "缺少 openpyxl，请安装 openpyxl 后重试，或另存为 CSV")
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        blocks: list[str] = []
        truncated = False
        for sheet in workbook.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue
            shown = rows[: row_limit + 1]
            if len(rows) > row_limit + 1:
                truncated = True
            width = max((len(row) for row in shown), default=0)
            lines = [f"[工作表：{sheet.title}]"]
            for row in shown:
                lines.append("| " + " | ".join("" if value is None else str(value) for value in row[:width]) + " |")
            blocks.append("\n".join(lines))
    except Exception as exc:
        return _failure("document_read_failed", f"XLSX 读取失败：{exc}")
    text = "\n\n".join(blocks)
    if not text:
        return _failure("empty_content", "XLSX 没有可读取内容")
    reason = f"工作表行数超过 {row_limit} 行，已截断。" if truncated else ""
    return _success(text, truncated=truncated, reason=reason)


def _shape_text(shape) -> list[str]:
    values: list[str] = []
    if getattr(shape, "has_text_frame", False):
        value = shape.text.strip()
        if value:
            values.append(value)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            values.append(" | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells))
    return values


def _parse_pptx(path: Path) -> ParserResult:
    try:
        from pptx import Presentation
    except (ImportError, ModuleNotFoundError):
        return _failure("dependency_missing", "缺少 python-pptx，请安装 python-pptx 后重试，或复制幻灯片文本")
    try:
        presentation = Presentation(str(path))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            values: list[str] = []
            for shape in slide.shapes:
                values.extend(_shape_text(shape))
            notes = getattr(slide, "notes_slide", None)
            if notes is not None:
                values.extend(
                    shape.text.strip()
                    for shape in notes.notes_text_frame.paragraphs
                    if shape.text.strip()
                )
            if values:
                slides.append(f"[幻灯片：{index}]\n" + "\n".join(values))
    except Exception as exc:
        return _failure("document_read_failed", f"PPTX 读取失败：{exc}")
    text = "\n\n".join(slides)
    return _success(text) if text else _failure("empty_content", "PPTX 没有可读取文本")


def parse_context_file(
    path: str | Path,
    *,
    max_bytes: int = DEFAULT_MAX_BYTES,
    csv_row_limit: int = DEFAULT_CSV_ROW_LIMIT,
) -> ParserResult:
    """Parse one supported file and never return failure text as正文."""
    path = Path(path)
    if not path.exists():
        return _failure("path_not_found", f"文件不存在：{path}")
    if not path.is_file():
        return _failure("path_not_file", f"路径不是文件：{path}")
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        hint = "请另存为 .docx/.xlsx/.pptx 或复制文本到剪贴板" if suffix in {".doc", ".xls", ".ppt"} else "当前版本不支持此文件格式"
        return _failure("unsupported_format", hint)
    try:
        if path.stat().st_size > max_bytes:
            return _failure("single_file_too_large", f"文件超过 {max_bytes:,} 字节保护线")
    except OSError as exc:
        return _failure("document_read_failed", f"无法读取文件信息：{exc}")
    if suffix == ".json":
        return _parse_json(path)
    if suffix == ".csv":
        return _parse_csv(path, csv_row_limit)
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".xlsx":
        return _parse_xlsx(path, csv_row_limit)
    if suffix == ".pptx":
        return _parse_pptx(path)
    try:
        text, encoding = read_text_with_fallback(path)
    except (OSError, ValueError, UnicodeError) as exc:
        return _failure("encoding_failed", str(exc))
    return _success(text, encoding=encoding) if text else _failure("empty_content", "文件没有可读取内容")
